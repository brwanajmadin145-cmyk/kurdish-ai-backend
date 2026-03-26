from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def create_title_slide(prs, content):
    """Simple title slide - just text"""
    title_slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(title_slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = content['title']
    
    if 'subtitle' in content:
        subtitle.text = content['subtitle']
    else:
        subtitle.text = ""
    
    # Simple formatting
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.name = 'Calibri'
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    
    if subtitle.text:
        subtitle.text_frame.paragraphs[0].font.size = Pt(20)
        subtitle.text_frame.paragraphs[0].font.name = 'Calibri'
        subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(89, 89, 89)
    
    return slide


def create_content_slide(prs, content, slide_number=1):
    """Simple content slide - just text"""
    bullet_slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(bullet_slide_layout)
    
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = content['title']
    
    # Format title
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.name = 'Calibri'
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    title.text_frame.paragraphs[0].font.bold = True
    
    # Add bullet points
    text_frame = body.text_frame
    text_frame.clear()  # Clear existing text
    
    points = content.get('points', [])
    
    for i, point in enumerate(points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = point
        p.level = 0
        p.font.size = Pt(18)
        p.font.name = 'Calibri'
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    return slide


def generate_gamma_presentation(topic, num_slides=5, theme='modern', ai_content=None):
    """
    Generate simple, clean PowerPoint presentation - just text, no colors
    
    Args:
        topic: Presentation topic
        num_slides: Number of slides
        theme: Not used (kept for compatibility)
        ai_content: AI-generated content
    
    Returns:
        Presentation object
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Prepare content
    if not ai_content or len(ai_content) == 0:
        ai_content = [
            {'title': topic, 'subtitle': 'Professional Presentation'},
        ] + [
            {
                'title': f'Section {i}',
                'points': ['Point 1', 'Point 2', 'Point 3', 'Point 4']
            }
            for i in range(1, num_slides)
        ]
    
    # Ensure correct number of slides
    while len(ai_content) < num_slides:
        ai_content.append({
            'title': 'Additional Points',
            'points': [
                'Key insight',
                'Important detail',
                'Essential information',
                'Critical factor'
            ]
        })
    
    ai_content = ai_content[:num_slides]
    
    # Create slides
    for i, content in enumerate(ai_content):
        if i == 0:
            # First slide is title
            create_title_slide(prs, content)
        else:
            # All other slides are content
            create_content_slide(prs, content, i)
    
    return prs


if __name__ == "__main__":
    # Test
    test_content = [
        {
            'title': 'Artificial Intelligence',
            'subtitle': 'The Future of Technology'
        },
        {
            'title': 'What is AI?',
            'points': [
                'Machine learning systems that learn from data',
                'Neural networks that mimic human brain',
                'Deep learning for complex pattern recognition',
                'Natural language processing for understanding text'
            ]
        },
        {
            'title': 'Applications',
            'points': [
                'Healthcare: Disease diagnosis and drug discovery',
                'Finance: Fraud detection and trading',
                'Education: Personalized learning platforms',
                'Transportation: Autonomous vehicles'
            ]
        },
        {
            'title': 'Key Takeaways',
            'points': [
                'AI is transforming every industry',
                'Continuous learning is essential',
                'Ethical considerations are important',
                'The future is AI-powered'
            ]
        }
    ]
    
    prs = generate_gamma_presentation('AI', 4, 'modern', test_content)
    prs.save("test_simple_text.pptx")
    print("✅ Simple text-only presentation created!")